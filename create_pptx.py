#!/usr/bin/env python3
"""Generate a stunning Calligro teacher recruitment PowerPoint presentation in Arabic."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
GOLD = RGBColor(0xD4, 0xAF, 0x37)
GOLD_LIGHT = RGBColor(0xFF, 0xEF, 0x96)
DARK = RGBColor(0x0A, 0x0A, 0x0A)
DARK2 = RGBColor(0x14, 0x14, 0x14)
DARK3 = RGBColor(0x1E, 0x1E, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE70 = RGBColor(0xB3, 0xB3, 0xB3)
WHITE50 = RGBColor(0x80, 0x80, 0x80)
RED = RGBColor(0xE7, 0x4C, 0x3C)
RED_DARK = RGBColor(0x2D, 0x0A, 0x06)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide):
    """Add dark background with a subtle gold accent shape."""
    set_bg(slide, DARK)
    # Gold accent circle (top-right)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x2A, 0x22, 0x0A)
    s.line.fill.background()
    # Bottom accent
    s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(5), Inches(5))
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x05)
    s2.line.fill.background()

def add_text(slide, left, top, width, height, text, size=24, bold=False, color=WHITE, align=PP_ALIGN.RIGHT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def add_gold_line(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()
    return s

def add_badge(slide, left, top, text, bg=GOLD, text_color=DARK):
    w, h = Inches(2.5), Inches(0.45)
    s = add_rounded_rect(slide, left, top, w, h, bg)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return s

def add_icon_card(slide, left, top, icon, title, desc, w=Inches(3.8), h=Inches(2.8)):
    card = add_rounded_rect(slide, left, top, w, h, RGBColor(0x1A, 0x1A, 0x1A), RGBColor(0x33, 0x33, 0x33))
    # Icon
    add_text(slide, left + Inches(0.3), top + Inches(0.25), Inches(1), Inches(0.7), icon, size=36, align=PP_ALIGN.RIGHT)
    # Title
    add_text(slide, left + Inches(0.3), top + Inches(0.9), w - Inches(0.6), Inches(0.5), title, size=18, bold=True, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    # Desc
    add_text(slide, left + Inches(0.3), top + Inches(1.45), w - Inches(0.6), Inches(1.2), desc, size=13, color=WHITE70, align=PP_ALIGN.RIGHT)

def add_stat(slide, left, top, number, label, num_color=GOLD):
    add_text(slide, left, top, Inches(2.5), Inches(0.8), number, size=44, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.75), Inches(2.5), Inches(0.5), label, size=13, color=WHITE50, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(sl)
# Decorative top gold bar
add_gold_line(sl, Inches(4.5), Inches(0), Inches(4.3))
# Badge
add_badge(sl, Inches(5.4), Inches(2.2), 'أكاديمية الخط العربي')
# Title
add_text(sl, Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.5), 'مرحباً بكم في كاليكرو', size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# Gold underline
add_gold_line(sl, Inches(5.5), Inches(4.4), Inches(2.3))
# Subtitle
add_text(sl, Inches(2), Inches(4.7), Inches(9.3), Inches(1), 'المنصة الرقمية الأولى لتعليم فن الخط العربي الأصيل', size=22, color=WHITE70, align=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(5.3), Inches(9.3), Inches(0.6), 'حيث يلتقي التراث بالتكنولوجيا', size=18, color=WHITE50, align=PP_ALIGN.CENTER)
# Bottom bar
add_gold_line(sl, Inches(4.5), Inches(7.35), Inches(4.3))

# ══════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(sl)
# Red accent
s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x2D, 0x0A, 0x06); s.line.fill.background()

add_badge(sl, Inches(5.4), Inches(0.6), 'التحدي', bg=RGBColor(0x3D, 0x15, 0x10), text_color=RED)
add_text(sl, Inches(1.5), Inches(1.3), Inches(10.3), Inches(1), 'واقع تعليم الخط العربي اليوم', size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_gold_line(sl, Inches(5.5), Inches(2.3), Inches(2.3))
add_text(sl, Inches(2), Inches(2.6), Inches(9.3), Inches(0.8), 'يواجه معلمو الخط العربي تحديات كبيرة في الوصول للطلاب وتنظيم دوراتهم وتحقيق دخل مستدام', size=16, color=WHITE70, align=PP_ALIGN.CENTER)

# Problem cards
problems = [
    ('❌', 'الوصول المحدود', 'التعليم مقتصر على مدينتك فقط\nلا يمكنك الوصول لطلاب العالم'),
    ('❌', 'بدون حضور رقمي', 'لا يوجد ملف تعريفي احترافي\nصعوبة التسويق لدوراتك'),
    ('❌', 'إدارة يدوية مرهقة', 'جدولة الحصص يدوياً\nتحصيل الرسوم بصعوبة'),
]
for i, (icon, title, desc) in enumerate(problems):
    x = Inches(0.8) + i * Inches(4.2)
    add_icon_card(sl, x, Inches(3.8), icon, title, desc, w=Inches(3.8), h=Inches(3))

# ══════════════════════════════════════════
# SLIDE 3: BEFORE vs AFTER (THE KEY SLIDE)
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK)

add_badge(sl, Inches(5.4), Inches(0.4), 'المقارنة')
add_text(sl, Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.8), 'قبل كاليكرو  ←→  بعد كاليكرو', size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# BEFORE card (right side for RTL)
before_card = add_rounded_rect(sl, Inches(6.9), Inches(2.1), Inches(5.8), Inches(5), RGBColor(0x1A, 0x0A, 0x08), RED)
add_text(sl, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.6), '❌  قبل كاليكرو', size=24, bold=True, color=RED, align=PP_ALIGN.RIGHT)
before_items = [
    'تعليم محلي محدود بمدينتك فقط',
    'لا يوجد ملف تعريفي احترافي على الإنترنت',
    'جدولة يدوية مرهقة للحصص والمواعيد',
    'صعوبة تحصيل الرسوم من الطلاب',
    'عدد طلاب محدود جداً في كل فصل',
    'لا توجد أدوات تفاعلية رقمية للتعليم',
]
for i, item in enumerate(before_items):
    y = Inches(3.05) + i * Inches(0.55)
    add_text(sl, Inches(7.4), y, Inches(5), Inches(0.5), f'•  {item}', size=14, color=RGBColor(0xCC, 0x88, 0x88), align=PP_ALIGN.RIGHT)

# AFTER card (left side for RTL)
after_card = add_rounded_rect(sl, Inches(0.6), Inches(2.1), Inches(5.8), Inches(5), RGBColor(0x14, 0x1A, 0x08), GOLD)
add_text(sl, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.6), '✅  بعد كاليكرو', size=24, bold=True, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
after_items = [
    'تعليم عالمي — طلاب من كل أنحاء العالم',
    'صفحة شخصية احترافية مميزة ومعرض أعمال',
    'جدولة تلقائية ذكية وإدارة متكاملة',
    'نظام دفع آمن عالمي (Apple Pay, بطاقات)',
    'عدد غير محدود من الطلاب في دوراتك',
    'فصول تفاعلية مباشرة بأحدث التقنيات',
]
for i, item in enumerate(after_items):
    y = Inches(3.05) + i * Inches(0.55)
    add_text(sl, Inches(0.8), y, Inches(5), Inches(0.5), f'•  {item}', size=14, color=RGBColor(0xBB, 0xCC, 0x88), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════
# SLIDE 4: WHAT WE OFFER
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(sl)

add_badge(sl, Inches(5.4), Inches(0.4), 'المنصة')
add_text(sl, Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.8), 'ماذا نوفر لك؟', size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

features = [
    ('🎓', 'لوحة تحكم متكاملة', 'إدارة الدورات والطلاب\nوالأرباح من مكان واحد'),
    ('📱', 'تطبيق جوال احترافي', 'تطبيق iOS و Android\nبتصميم عالمي مميز'),
    ('🌐', 'بوابة ويب متطورة', 'موقع إلكتروني احترافي\nيعرض دوراتك للعالم'),
    ('💳', 'نظام دفع عالمي', 'Apple Pay و Google Pay\nوبطاقات الائتمان'),
    ('🔔', 'إشعارات ذكية', 'تنبيهات فورية للطلاب\nبالمواعيد والتحديثات'),
    ('🏆', 'معرض أعمالك', 'عرض إبداعاتك في معرض\nرقمي احترافي متكامل'),
]
for i, (icon, title, desc) in enumerate(features):
    row = i // 3
    col = i % 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(2.2) + row * Inches(2.5)
    add_icon_card(sl, x, y, icon, title, desc, w=Inches(3.8), h=Inches(2.2))

# ══════════════════════════════════════════
# SLIDE 5: HOW YOU EARN
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(sl)

add_badge(sl, Inches(5.4), Inches(0.5), 'الأرباح')
add_text(sl, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8), 'كيف تربح مع كاليكرو؟', size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(2.1), Inches(9.3), Inches(0.6), 'نظام شفاف وعادل يضمن لك أعلى عائد من دوراتك', size=16, color=WHITE70, align=PP_ALIGN.CENTER)

# Stats
add_stat(sl, Inches(1), Inches(3.0), '٥٠٪', 'خصم للطلاب عبر الموقع')
add_stat(sl, Inches(4), Inches(3.0), '∞', 'عدد غير محدود من الطلاب')
add_stat(sl, Inches(7), Inches(3.0), 'عالمي', 'وصول لكل دول العالم')
add_stat(sl, Inches(10), Inches(3.0), 'فوري', 'تحويل الأرباح مباشرة')

# Example box
ex = add_rounded_rect(sl, Inches(2.5), Inches(4.8), Inches(8.3), Inches(2.2), RGBColor(0x1A, 0x18, 0x08), GOLD)
add_text(sl, Inches(3), Inches(5.0), Inches(7.3), Inches(0.5), '💰 مثال عملي على الأرباح', size=20, bold=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, Inches(3), Inches(5.6), Inches(7.3), Inches(0.5), 'دورة بسعر ١٠٠$ × ٣٠ طالب = ٣,٠٠٠$ لدورة واحدة!', size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(3), Inches(6.2), Inches(7.3), Inches(0.5), 'تخيل ٤ دورات في السنة = ١٢,٠٠٠$ دخل إضافي من فنك', size=16, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 6: VISION
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(sl)

add_badge(sl, Inches(5.4), Inches(0.5), 'الرؤية')
add_text(sl, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8), 'رؤيتنا المستقبلية', size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(2.1), Inches(9.3), Inches(0.8), 'أن نكون المنصة الأولى عالمياً\nفي تعليم فن الخط العربي الأصيل', size=22, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

vision = [
    ('🌍', 'انتشار عالمي', 'نهدف للوصول لملايين الطلاب\nحول العالم العربي والإسلامي'),
    ('🤝', 'مجتمع متكامل', 'بناء أكبر مجتمع رقمي\nلعشاق فن الخط العربي'),
    ('🏅', 'شهادات معتمدة', 'شهادات إتمام معتمدة\nمن أكاديمية كاليكرو'),
]
for i, (icon, title, desc) in enumerate(vision):
    x = Inches(0.8) + i * Inches(4.2)
    add_icon_card(sl, x, Inches(3.5), icon, title, desc, w=Inches(3.8), h=Inches(2.5))

# Bottom quote
add_text(sl, Inches(2), Inches(6.4), Inches(9.3), Inches(0.6), '« فنّك يستحق أن يصل للعالم »', size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 7: CTA - JOIN US
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK)
# Big gold glow
s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(1), Inches(7), Inches(7))
s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x05); s.line.fill.background()

add_text(sl, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2), 'انضم إلى كاليكرو اليوم', size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_gold_line(sl, Inches(5), Inches(3.2), Inches(3.3))
add_text(sl, Inches(2), Inches(3.5), Inches(9.3), Inches(0.8), 'كن جزءاً من مستقبل تعليم الخط العربي', size=24, color=WHITE70, align=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(4.2), Inches(9.3), Inches(0.6), '🌟 فنّك يستحق أن يصل للعالم 🌟', size=20, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# CTA Button
btn = add_rounded_rect(sl, Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.9), GOLD)
tf = btn.text_frame
p = tf.paragraphs[0]
p.text = 'ابدأ رحلتك الآن'
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK
p.font.name = 'Arial'
p.alignment = PP_ALIGN.CENTER

add_text(sl, Inches(2), Inches(6.5), Inches(9.3), Inches(0.5), 'www.calligro.com', size=14, color=WHITE50, align=PP_ALIGN.CENTER)

# Save
output = os.path.expanduser('~/Desktop/Calligro_Teacher_Presentation.pptx')
prs.save(output)
print(f'✅ Presentation saved to: {output}')
