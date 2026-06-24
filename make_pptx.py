from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

G=RGBColor(0xD4,0xAF,0x37);GL=RGBColor(0xFF,0xEF,0x96);D=RGBColor(0x0A,0x0A,0x0A)
W=RGBColor(0xFF,0xFF,0xFF);W7=RGBColor(0xB0,0xB0,0xB0);W4=RGBColor(0x70,0x70,0x70)
R=RGBColor(0xE7,0x4C,0x3C);DG=RGBColor(0x14,0x14,0x14);D2=RGBColor(0x1A,0x1A,0x1A)

prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5)

def bg(sl,c=D):
    f=sl.background.fill;f.solid();f.fore_color.rgb=c

def rect(sl,l,t,w,h,fc,lc=None):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    s.fill.solid();s.fill.fore_color.rgb=fc
    if lc:s.line.color.rgb=lc;s.line.width=Pt(1.5)
    else:s.line.fill.background()
    return s

def txt(sl,l,t,w,h,s,sz=18,bold=False,col=W,al=PP_ALIGN.RIGHT):
    tb=sl.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=s;p.font.size=Pt(sz);p.font.bold=bold
    p.font.color.rgb=col;p.font.name='Arial';p.alignment=al
    return tb

def line(sl,l,t,w,col=G):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,Inches(0.04))
    s.fill.solid();s.fill.fore_color.rgb=col;s.line.fill.background()

def oval(sl,l,t,w,h,fc):
    s=sl.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    s.fill.solid();s.fill.fore_color.rgb=fc;s.line.fill.background()

def badge(sl,l,t,s,bg2=G,tc=D):
    r=rect(sl,l,t,Inches(2.7),Inches(0.45),bg2)
    tf=r.text_frame;p=tf.paragraphs[0];p.text=s
    p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=tc
    p.font.name='Arial';p.alignment=PP_ALIGN.CENTER

def card(sl,l,t,w,h,icon,title,desc):
    rect(sl,l,t,w,h,D2,RGBColor(0x33,0x33,0x33))
    txt(sl,l+Inches(.2),t+Inches(.2),w-Inches(.4),Inches(.55),icon,34,al=PP_ALIGN.RIGHT)
    txt(sl,l+Inches(.2),t+Inches(.82),w-Inches(.4),Inches(.4),title,16,True,GL,PP_ALIGN.RIGHT)
    txt(sl,l+Inches(.2),t+Inches(1.28),w-Inches(.4),h-Inches(1.4),desc,12,False,W7,PP_ALIGN.RIGHT)

def stat(sl,l,t,num,lbl,nc=G):
    txt(sl,l,t,Inches(2.8),Inches(.8),num,42,True,nc,PP_ALIGN.CENTER)
    txt(sl,l,t+Inches(.8),Inches(2.8),Inches(.45),lbl,12,False,W4,PP_ALIGN.CENTER)

# SLIDE 1
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
oval(s,Inches(8.5),Inches(-1.5),Inches(6),Inches(6),RGBColor(0x20,0x18,0x05))
oval(s,Inches(-2),Inches(4),Inches(5),Inches(5),RGBColor(0x18,0x12,0x04))
line(s,Inches(4.5),Inches(0.08),Inches(4.3))
badge(s,Inches(5.3),Inches(1.8),'⭐  أكاديمية الخط العربي الرقمية')
txt(s,Inches(1),Inches(2.5),Inches(11.3),Inches(1.4),'مرحباً بكم في كاليكرو',56,True,W,PP_ALIGN.CENTER)
line(s,Inches(4.7),Inches(3.95),Inches(3.9))
txt(s,Inches(1.5),Inches(4.2),Inches(10.3),Inches(.6),'المنصة الرقمية الأولى لتعليم فن الخط العربي الأصيل',22,False,W7,PP_ALIGN.CENTER)
txt(s,Inches(1.5),Inches(4.9),Inches(10.3),Inches(.5),'حيث يلتقي التراث بالتكنولوجيا الحديثة',16,False,W4,PP_ALIGN.CENTER)
r=rect(s,Inches(4.5),Inches(5.8),Inches(4.3),Inches(.8),G)
tf=r.text_frame;p=tf.paragraphs[0];p.text='نبدأ الرحلة ←'
p.font.size=Pt(20);p.font.bold=True;p.font.color.rgb=D;p.font.name='Arial';p.alignment=PP_ALIGN.CENTER
line(s,Inches(4.5),Inches(7.38),Inches(4.3))

# SLIDE 2
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
oval(s,Inches(9),Inches(-2),Inches(5),Inches(5),RGBColor(0x2A,0x08,0x05))
badge(s,Inches(5.3),Inches(.4),'❗  التحديات الحالية',RGBColor(0x3D,0x10,0x08),R)
txt(s,Inches(1),Inches(1.1),Inches(11.3),Inches(.8),'واقع معلمي الخط العربي اليوم',40,True,W,PP_ALIGN.CENTER)
line(s,Inches(5.2),Inches(2.0),Inches(2.9),R)
txt(s,Inches(1.5),Inches(2.3),Inches(10.3),Inches(.5),'الكثير من المواهب تضيع بسبب غياب المنصة المناسبة',16,False,W7,PP_ALIGN.CENTER)
probs=[('❌','وصول محدود','تعليمك محصور في مدينتك\nلا يصل لطلاب العالم'),
       ('❌','بلا هوية رقمية','لا ملف احترافي على الإنترنت\nصعوبة التسويق لنفسك'),
       ('❌','إدارة مرهقة','جداول يدوية وتحصيل صعب\nيضيع وقتك الثمين')]
for i,(ic,ti,de) in enumerate(probs):
    card(s,Inches(.5+i*4.3),Inches(3.1),Inches(4.0),Inches(2.9),ic,ti,de)

# SLIDE 3
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
badge(s,Inches(5.3),Inches(.35),'🔄  التحول')
txt(s,Inches(1),Inches(1.0),Inches(11.3),Inches(.7),'قبل كاليكرو  ↔  بعد كاليكرو',38,True,W,PP_ALIGN.CENTER)
rect(s,Inches(6.9),Inches(1.7),Inches(5.9),Inches(5.5),RGBColor(0x1C,0x08,0x06),R)
txt(s,Inches(7.2),Inches(1.85),Inches(5.4),Inches(.5),'❌  قبل كاليكرو',22,True,R,PP_ALIGN.RIGHT)
for i,t in enumerate(['تعليم محلي — محدود بمدينتك','بلا هوية رقمية احترافية','جداول وإدارة يدوية','صعوبة تحصيل الرسوم','عدد طلاب محدود جداً','لا أدوات تعليمية تفاعلية']):
    txt(s,Inches(7.4),Inches(2.55)+Inches(i*.52),Inches(5.3),Inches(.45),f'•  {t}',13.5,False,RGBColor(0xCC,0x88,0x88),PP_ALIGN.RIGHT)
rect(s,Inches(.5),Inches(1.7),Inches(5.9),Inches(5.5),RGBColor(0x10,0x18,0x05),G)
txt(s,Inches(.7),Inches(1.85),Inches(5.4),Inches(.5),'✅  بعد كاليكرو',22,True,GL,PP_ALIGN.RIGHT)
for i,t in enumerate(['وصول عالمي — طلاب من كل مكان','صفحة احترافية ومعرض أعمال','جدولة تلقائية وإدارة متكاملة','نظام دفع آمن عالمي','عدد غير محدود من الطلاب','فصول تفاعلية بأحدث التقنيات']):
    txt(s,Inches(.7),Inches(2.55)+Inches(i*.52),Inches(5.3),Inches(.45),f'•  {t}',13.5,False,RGBColor(0xAA,0xCC,0x77),PP_ALIGN.RIGHT)

# SLIDE 4
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
oval(s,Inches(5),Inches(-2),Inches(6),Inches(6),RGBColor(0x1A,0x15,0x04))
badge(s,Inches(5.3),Inches(.35),'🚀  المنصة')
txt(s,Inches(1),Inches(1.0),Inches(11.3),Inches(.7),'ماذا نوفر لك؟',42,True,W,PP_ALIGN.CENTER)
txt(s,Inches(1.5),Inches(1.8),Inches(10.3),Inches(.5),'كل ما تحتاجه في مكان واحد',16,False,W7,PP_ALIGN.CENTER)
feats=[('🎓','لوحة تحكم متكاملة','إدارة دوراتك وطلابك وأرباحك\nبسهولة من مكان واحد'),
       ('📱','تطبيق iOS و Android','تطبيق احترافي بتصميم عالمي\nللمعلمين والطلاب'),
       ('🌐','بوابة ويب احترافية','موقع يعرض دوراتك للعالم\nعلى مدار الساعة'),
       ('💳','دفع عالمي آمن','Apple Pay, Google Pay\nوبطاقات الائتمان'),
       ('🔔','إشعارات فورية','تنبيهات تلقائية للطلاب\nبالمواعيد والتحديثات'),
       ('🏆','معرض أعمال رقمي','اعرض إبداعاتك في معرض\naحترافي مميز')]
for i,(ic,ti,de) in enumerate(feats):
    r2=i//3;c=i%3
    card(s,Inches(.4+c*4.3),Inches(2.55+r2*2.3),Inches(4.0),Inches(2.05),ic,ti,de)

# SLIDE 5
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
oval(s,Inches(4),Inches(1),Inches(5),Inches(5),RGBColor(0x20,0x18,0x05))
badge(s,Inches(5.3),Inches(.35),'💰  الأرباح')
txt(s,Inches(1),Inches(1.0),Inches(11.3),Inches(.7),'كيف تربح مع كاليكرو؟',42,True,W,PP_ALIGN.CENTER)
txt(s,Inches(1.5),Inches(1.8),Inches(10.3),Inches(.5),'نظام شفاف يضمن لك أعلى عائد من فنك',16,False,W7,PP_ALIGN.CENTER)
stat(s,Inches(.3),Inches(2.8),'٥٠٪','خصم ترحيبي للطلاب')
stat(s,Inches(3.5),Inches(2.8),'∞','طلاب بلا حد')
stat(s,Inches(6.7),Inches(2.8),'🌍','وصول عالمي')
stat(s,Inches(9.9),Inches(2.8),'⚡','دفع فوري')
r=rect(s,Inches(1.8),Inches(4.6),Inches(9.7),Inches(2.5),RGBColor(0x18,0x15,0x05),G)
txt(s,Inches(2.2),Inches(4.85),Inches(9),Inches(.5),'💰 مثال عملي — اكتشف دخلك المحتمل',20,True,GL,PP_ALIGN.CENTER)
line(s,Inches(4),Inches(5.45),Inches(5.3))
txt(s,Inches(2.2),Inches(5.6),Inches(9),Inches(.45),'دورة بسعر ١٠٠$ × ٣٠ طالب = 3,000$ لدورة واحدة!',22,True,W,PP_ALIGN.CENTER)
txt(s,Inches(2.2),Inches(6.15),Inches(9),Inches(.5),'٤ دورات سنوياً = دخل يصل إلى 12,000$ من فنك',15,False,G,PP_ALIGN.CENTER)

# SLIDE 6
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
badge(s,Inches(5.3),Inches(.35),'📋  الخطوات')
txt(s,Inches(1),Inches(1.0),Inches(11.3),Inches(.7),'كيف تبدأ في ٣ خطوات فقط؟',42,True,W,PP_ALIGN.CENTER)
steps=[('①','إنشاء حسابك','سجّل كمعلم واحصل على\nملفك الاحترافي\nفي دقائق معدودة',RGBColor(0x14,0x14,0x20)),
       ('②','أنشئ دورتك','حدد المحتوى، الجدول\nوالسعر وعدد الطلاب\nبسهولة تامة',RGBColor(0x14,0x18,0x08)),
       ('③','استقبل طلابك','انشر دورتك وابدأ\nالتدريس واستلم\nأرباحك فوراً',RGBColor(0x10,0x18,0x05))]
for i,(num,ti,de,fc2) in enumerate(steps):
    rect(s,Inches(.5+i*4.3),Inches(2.2),Inches(4.0),Inches(4.5),fc2,G)
    txt(s,Inches(.7+i*4.3),Inches(2.4),Inches(3.6),Inches(.8),num,52,True,G,PP_ALIGN.CENTER)
    txt(s,Inches(.7+i*4.3),Inches(3.3),Inches(3.6),Inches(.5),ti,20,True,GL,PP_ALIGN.CENTER)
    txt(s,Inches(.7+i*4.3),Inches(3.95),Inches(3.6),Inches(2.0),de,15,False,W7,PP_ALIGN.CENTER)
    if i<2:
        txt(s,Inches(4.4+i*4.3),Inches(4.1),Inches(.5),Inches(.5),'←',26,True,G,PP_ALIGN.CENTER)

# SLIDE 7
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
oval(s,Inches(4.5),Inches(0),Inches(6),Inches(6),RGBColor(0x1A,0x14,0x04))
badge(s,Inches(5.3),Inches(.35),'🌟  الرؤية')
txt(s,Inches(1),Inches(1.0),Inches(11.3),Inches(.7),'رؤيتنا المستقبلية',42,True,W,PP_ALIGN.CENTER)
txt(s,Inches(1.5),Inches(1.8),Inches(10.3),Inches(.6),'أن نكون المنصة الأولى عالمياً لتعليم فن الخط العربي',20,False,GL,PP_ALIGN.CENTER)
vis=[('🌍','انتشار عالمي','ملايين الطلاب حول العالم\nالعربي والإسلامي'),
     ('🤝','مجتمع متكامل','أكبر مجتمع رقمي\nلعشاق الخط العربي'),
     ('🏅','شهادات معتمدة','شهادات إتمام دولية\nمن أكاديمية كاليكرو')]
for i,(ic,ti,de) in enumerate(vis):
    card(s,Inches(.5+i*4.3),Inches(2.9),Inches(4.0),Inches(3.0),ic,ti,de)
line(s,Inches(3.5),Inches(6.5),Inches(6.3))
txt(s,Inches(1.5),Inches(6.6),Inches(10.3),Inches(.55),'« فنّك يستحق أن يصل للعالم »',18,True,G,PP_ALIGN.CENTER)

# SLIDE 8
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
oval(s,Inches(2.5),Inches(.5),Inches(8),Inches(8),RGBColor(0x1E,0x16,0x04))
line(s,Inches(4.5),Inches(.08),Inches(4.3))
txt(s,Inches(1),Inches(1.5),Inches(11.3),Inches(1.2),'انضم إلى كاليكرو اليوم',52,True,W,PP_ALIGN.CENTER)
line(s,Inches(4.5),Inches(2.85),Inches(4.3))
txt(s,Inches(1.5),Inches(3.1),Inches(10.3),Inches(.6),'كن جزءاً من مستقبل تعليم الخط العربي',24,False,W7,PP_ALIGN.CENTER)
txt(s,Inches(1.5),Inches(3.8),Inches(10.3),Inches(.6),'🌟  فنّك يستحق أن يصل للعالم  🌟',20,False,GL,PP_ALIGN.CENTER)
bens=['✦  ملف احترافي مميز يعرضك للعالم','✦  نظام دفع آمن وعالمي متكامل','✦  طلاب من كل أنحاء العالم','✦  دعم فني متواصل وتطوير مستمر']
for i,b in enumerate(bens):
    c=i%2;r2=i//2
    txt(s,Inches(.8+c*6.5),Inches(4.6+r2*.55),Inches(5.8),Inches(.45),b,15,False,W7,PP_ALIGN.RIGHT)
r=rect(s,Inches(4.2),Inches(6.0),Inches(4.9),Inches(.9),G)
tf=r.text_frame;p=tf.paragraphs[0];p.text='ابدأ رحلتك الآن ←'
p.font.size=Pt(22);p.font.bold=True;p.font.color.rgb=D;p.font.name='Arial';p.alignment=PP_ALIGN.CENTER
txt(s,Inches(1.5),Inches(7.05),Inches(10.3),Inches(.4),'www.calligro.com  |  أكاديمية الخط العربي',12,False,W4,PP_ALIGN.CENTER)

prs.save('/Users/yazanqattous/Desktop/Calligro_Teachers_Presentation.pptx')
print('Done!')
